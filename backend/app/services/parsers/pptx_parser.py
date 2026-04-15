"""PPTX 解析器。"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.services.parsers.base import BaseParser, ParsedBlock, ParsedDocument, normalize_text


class PptxParser(BaseParser):
    """提取每页标题、正文和表格。"""

    parser_name = "python-pptx"

    def parse(self, file_path: Path) -> ParsedDocument:
        presentation = Presentation(str(file_path))
        blocks: list[ParsedBlock] = []
        block_order = 0
        title = file_path.stem

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_title = ""
            if slide.shapes.title is not None:
                slide_title = slide.shapes.title.text.strip()
            if slide_title:
                if slide_index == 1:
                    title = slide_title
                blocks.append(
                    ParsedBlock(
                        block_type="title",
                        raw_text=slide_title,
                        normalized_text=normalize_text(slide_title),
                        section_path=[slide_title],
                        page_no=slide_index,
                        block_order=block_order,
                        metadata={"slide_no": slide_index},
                    )
                )
                block_order += 1

            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        values = [" ".join(cell.text.split()) for cell in row.cells]
                        rows.append(" | ".join(values))
                    table_text = "\n".join(row for row in rows if row.strip())
                    if table_text:
                        blocks.append(
                            ParsedBlock(
                                block_type="table",
                                raw_text=table_text,
                                normalized_text=normalize_text(table_text),
                                section_path=[slide_title] if slide_title else [],
                                page_no=slide_index,
                                block_order=block_order,
                                metadata={"slide_no": slide_index},
                            )
                        )
                        block_order += 1
                    continue

                if not getattr(shape, "has_text_frame", False):
                    continue
                text = "\n".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()).strip()
                if not text or (slide_title and text == slide_title):
                    continue

                blocks.append(
                    ParsedBlock(
                        block_type="list" if "\n" in text else "paragraph",
                        raw_text=text,
                        normalized_text=normalize_text(text),
                        section_path=[slide_title] if slide_title else [],
                        page_no=slide_index,
                        block_order=block_order,
                        metadata={"slide_no": slide_index},
                    )
                )
                block_order += 1

        return ParsedDocument(
            title=title,
            file_type="pptx",
            parser_name=self.parser_name,
            blocks=blocks,
            metadata={"page_count": len(presentation.slides)},
        )
